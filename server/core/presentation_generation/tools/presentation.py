from PIL import Image
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import random

from constants import RESPONSE_PATH
from custom_types.presentation_generation import SlideInput
def resize_and_crop_image(image_bytes, target_height=720, crop_width=480):
 
    # Load the image from BytesIO
    image = Image.open(image_bytes)

    # Calculate the new width to maintain the aspect ratio
    aspect_ratio = image.width / image.height
    target_width = int(target_height * aspect_ratio)

    # Resize the image
    image = image.resize((target_width, target_height), Image.LANCZOS)

    # Calculate cropping box
    left = (target_width - crop_width) // 2
    right = left + crop_width
    top = 0
    bottom = target_height

    # Crop the image
    cropped_image = image.crop((left, top, right, bottom))

    cropped_image_bytes = BytesIO()
    cropped_image.save(cropped_image_bytes, format="PNG")
    cropped_image_bytes.seek(0)

    return cropped_image_bytes
# Helper function to add a colored background shape

def add_image_from_bytes(slide, image_bytes, left, top, width=None, height=None):

    # Add the image to the slide
    image = slide.shapes.add_picture(image_bytes, left, top, width, height)

    return image

def add_colored_background(slide, presentation, color=(255, 255, 255)):
    """
    Adds a rectangle covering the entire slide, filled with the specified RGB color.
    This helps create a custom background color.
    """
    # Get slide dimensions from the presentation object
    width = presentation.slide_width
    height = presentation.slide_height
    
    # Add a rectangle shape to cover the entire slide
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, width, height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(*color)
    background.line.fill.background()  # No outline
    background.z_order = 0  # Send to back


# Helper function to style title text

def style_title(shape, font_name="Calibri", font_size=40, font_color=(0, 0, 0)):
    """
    Styles a given title shape with the specified font name, size, and RGB font color.
    """
    paragraphs = shape.text_frame.paragraphs
    for p in paragraphs:
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(*font_color)
        p.alignment = PP_ALIGN.LEFT

# Helper function to style body text

def style_body(shape, font_name="Calibri", font_size=20, font_color=(50, 50, 50)):
    """
    Styles the text within a shape's text frame.
    """
    for paragraph in shape.text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = RGBColor(*font_color)
        paragraph.alignment = PP_ALIGN.LEFT





def make_presentation(slide_input: SlideInput):
    """
    Generates a PowerPoint presentation based on the provided input data and template.

    This function creates a multi-slide presentation using a specified template and dynamically 
    inserts text as specified in the slide input.

    Args:
        slide_input (SlideInput): Configuration object containing:
            - slide_data (SlideData): Data for each slide in the presentation
            - template_name (str): The name of the template folder to be used
            - id (str): A unique identifier for the presentation file
    """

    slide_data = slide_input.slide_data
    template_name = slide_input.template_name
    id = slide_input.id

    presentation = Presentation()
 
    assets_img = [f'pptassets/pptimg{i}.jpeg' for i in range(1, 20)]

    # Randomly choose 10 unique images from the array
    selected_images = random.sample(assets_img, 10)

    ############################
    # Slide 1 - Startup Title and Slogan
    ############################
    slide1 = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide1.shapes.add_picture(
        f'{template_name}/1.png', Inches(0), Inches(0), width=presentation.slide_width, height=presentation.slide_height
    )

    # Title Text
    title1 = slide1.shapes.add_textbox(
        Inches(0), Inches(2), presentation.slide_width, Inches(1)
    )
    text_frame1 = title1.text_frame
    text_frame1.text = slide_data.startup_title
    text_frame1.paragraphs[0].font.size = Pt(40)
    text_frame1.paragraphs[0].font.bold = True
    text_frame1.paragraphs[0].font.name = "Arial"
    text_frame1.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame1.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # Slogan Text
    slogan = slide1.shapes.add_textbox(
        Inches(0), Inches(3), presentation.slide_width, Inches(1)
    )
    slogan_frame = slogan.text_frame
    slogan_frame.text = slide_data.slogan
    slogan_frame.paragraphs[0].font.size = Pt(20)
    slogan_frame.paragraphs[0].font.bold = False
    slogan_frame.paragraphs[0].font.name = "Arial"
    slogan_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    slogan_frame.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)

    ############################
    # Slide 2 - About Us
    ############################
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide2.shapes.add_picture(
        f'{template_name}/2.png', Inches(0), Inches(0), width=presentation.slide_width, height=presentation.slide_height
    )

    # Left Image
    slide2.shapes.add_picture(
        f'{selected_images[0]}', Inches(0.5), Inches(2), width=Inches(4.25), height=Inches(2.5)
    )

    # Title Text
    title = slide2.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(1))
    title_frame = title.text_frame
    title_frame.text = slide_data.about_us_title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = "Arial"
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # Body Text
    body = slide2.shapes.add_textbox(Inches(5), Inches(3), Inches(4), Inches(3))
    body_frame = body.text_frame
    body_frame.text = slide_data.about_us_text
    for paragraph in body_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.name = "Arial"
        paragraph.font.color.rgb = RGBColor(80, 80, 80)

    ############################
    # Slide 3 - Main Goals
    ############################
    slide3 = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide3.shapes.add_picture(
        f'{template_name}/3.png', Inches(0), Inches(0), width=presentation.slide_width, height=presentation.slide_height
    )

    # Title Text
    title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title.text_frame
    title_frame.text = slide_data.main_goals_title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = "Arial"
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Body Text
    body = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7), Inches(2))
    body_frame = body.text_frame
    body_frame.text = slide_data.main_goals_text
    for paragraph in body_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.name = "Arial"
        paragraph.font.color.rgb = RGBColor(80, 80, 80)
        paragraph.alignment = PP_ALIGN.LEFT

    # Add Image
    with open(f'{selected_images[1]}', "rb") as f:
        image_bytes1 = BytesIO(f.read())

    with open(f'{selected_images[2]}', "rb") as f:
        image_bytes2 = BytesIO(f.read())

    slide3.shapes.add_picture(image_bytes1, Inches(0.5), Inches(4.25), Inches(4.25), Inches(2.5))
    slide3.shapes.add_picture(image_bytes2, Inches(5), Inches(4.25), Inches(4.25), Inches(2.5))

    ############################
    # Slide 4 - Problems and Solutions
    ############################
    slide4 = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide4.shapes.add_picture(
        f'{template_name}/4.png', Inches(0), Inches(0), width=presentation.slide_width, height=presentation.slide_height
    )

    title = slide4.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
    title_frame = title.text_frame
    title_frame.text = slide_data.slide4.title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = "Arial"
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    body = slide4.shapes.add_textbox(Inches(1), Inches(3), Inches(7), Inches(2))
    body_frame = body.text_frame
    body_frame.text = slide_data.slide4.body
    for paragraph in body_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.name = "Arial"
        paragraph.font.color.rgb = RGBColor(80, 80, 80)

    ############################
    # Slide 5 - Product
    ############################
    slide5 = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide5.shapes.add_picture(
        f'{template_name}/5.png', Inches(0), Inches(0), width=presentation.slide_width, height=presentation.slide_height
    )

    title = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.25), Inches(2.5))
    title_frame = title.text_frame
    title_frame.text = slide_data.slide5.title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = "Arial"
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    body = slide5.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(4.25), Inches(2.5))
    body_frame = body.text_frame
    body_frame.text = slide_data.slide5.body
    for paragraph in body_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.name = "Arial"
        paragraph.font.color.rgb = RGBColor(80, 80, 80)

    with open(f'{selected_images[3]}', "rb") as f:
        slide5.shapes.add_picture(
            BytesIO(f.read()), Inches(5), Inches(2.5), Inches(4), Inches(3)
        )

    ############################
    # Slide 6 - Business/Revenue Model
    ############################
    slide6 = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide6.shapes.add_picture(
        f'{template_name}/6.png', Inches(0), Inches(0), width=presentation.slide_width, height=presentation.slide_height
    )

    title = slide6.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1.5))
    title_frame = title.text_frame
    title_frame.text = slide_data.slide6.title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = "Arial"
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    body = slide6.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(3))
    body_frame = body.text_frame
    body_frame.text = slide_data.slide6.body
    for paragraph in body_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.name = "Arial"
        paragraph.font.color.rgb = RGBColor(80, 80, 80)

    y_position = 0.6
    image_grid_1 = selected_images[4]
    image_grid_2 =selected_images[5]
    image_grid_3 =selected_images[6]
    image_grid_4 = selected_images[7]

    for i, service in enumerate(slide_data.slide6.services):
        # Use i to index the `selected_images` list
        slide6.shapes.add_picture(selected_images[i + 4], Inches(5), Inches(y_position), Inches(2.125), Inches(1.25))

        # Add the service title and description
        text_box = slide6.shapes.add_textbox(Inches(7.2), Inches(y_position), Inches(2.7), Inches(0.8))
        text_frame = text_box.text_frame

        # Service title
        p = text_frame.add_paragraph()
        p.text = service.title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = "Arial"
        p.font.color.rgb = RGBColor(0, 0, 0)

        # Service description
        p = text_frame.add_paragraph()
        p.text = service.description
        p.font.size = Pt(14)
        p.font.name = "Arial"
        p.font.color.rgb = RGBColor(80, 80, 80)

        # Increment the vertical position
        y_position += 1.51


    ############################
    # Slide 7 - Status and Milestones
    ############################
    slide7 = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide7.shapes.add_picture(
        f'{template_name}/7.png', Inches(0), Inches(0), width=presentation.slide_width, height=presentation.slide_height
    )

    slide7.shapes.add_picture(selected_images[8], Inches(0.5), Inches(0.5), Inches(4.25), Inches(2.5))
    slide7.shapes.add_picture(selected_images[9], Inches(5), Inches(0.5), Inches(4.25), Inches(2.5))

    title = slide7.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    title_frame = title.text_frame
    title_frame.text = slide_data.slide7.title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = "Arial"
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    body = slide7.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
    body_frame = body.text_frame
    body_frame.text = slide_data.slide7.body
    for paragraph in body_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.name = "Arial"
        paragraph.font.color.rgb = RGBColor(80, 80, 80)


    ############################
    # Slide 8 - Marketing and Sales Strategy
    ############################


    # Add a new slide with blank layout
    slide8 = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide8.shapes.add_picture(
            f'{template_name}/8.png', Inches(0), Inches(0), width=presentation.slide_width, height=presentation.slide_height
        )
    # Add the "THANK YOU" title
    title = slide8.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_frame = title.text_frame
    title_frame.text = "THANK YOU"
    title_frame.paragraphs[0].font.size = Pt(48)  # Large font for the title
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = "Arial"
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black color
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center alignment

    # Add the "FOR YOUR ATTENTION" subtitle
    subtitle = slide8.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = "FOR YOUR ATTENTION"
    subtitle_frame.paragraphs[0].font.size = Pt(24)  # Medium font size for the subtitle
    subtitle_frame.paragraphs[0].font.bold = False
    subtitle_frame.paragraphs[0].font.name = "Arial"
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)  # Gray color
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center alignment



    presentation.save(f"{RESPONSE_PATH}/presentation_{id}.pptx")


    print("Presentation created successfully and saved as '{RESPONSE_PATH}/presentation_{id}.pptx'")

    return {
        "message": "Presentation created successfully.",
        "file_name": f"presentation_{id}.pptx"
    }




# print("Presentation created successfully and saved as 'Pet_Sitting_Presentation.pptx'")


# make_presentation(slide_data, "")